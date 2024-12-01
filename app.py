import streamlit as st
import tempfile
from pathlib import Path
from pptx import Presentation
from spellchecker import SpellChecker
import csv
import re
import string
import requests

# LanguageTool public API endpoint (no API key required)
LANGUAGE_TOOL_API = "https://api.languagetool.org/v2/check"

# Function to check grammar using LanguageTool API (no API Key required)
def check_grammar(text):
    if not text.strip():
        return "No issues", text  # Return "No issues" for empty text
    try:
        payload = {
            'text': text,
            'language': 'en-US',
            'enabledOnly': 'false'
        }
        response = requests.post(LANGUAGE_TOOL_API, data=payload)
        response_data = response.json()

        if 'matches' in response_data and response_data['matches']:
            # Apply suggested fixes
            corrected_text = text
            for match in reversed(response_data['matches']):  # Reverse to avoid index conflicts
                if 'replacement' in match and match['replacement']:
                    start = match['offset']
                    end = start + match['length']
                    corrected_text = corrected_text[:start] + match['replacement'][0] + corrected_text[end:]
            
            if corrected_text.lower() != text.lower():  # Only record if there are significant changes
                return "Grammatical error", corrected_text
        return "No issues", text
    except Exception as e:
        st.error(f"Error in grammar check: {e}")
        return "Error", text

# Function to validate grammar in a presentation
def validate_grammar(input_ppt):
    presentation = Presentation(input_ppt)
    grammar_issues = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip():
                            issue, corrected_text = check_grammar(run.text)
                            if issue != "No issues":
                                grammar_issues.append({
                                    'slide': slide_index,
                                    'issue': issue,
                                    'text': run.text,
                                    'corrected': corrected_text
                                })

    return grammar_issues

# Function to detect and correct misspellings
def detect_misspellings(text):
    spell = SpellChecker()
    words = text.split()
    misspellings = {}

    for word in words:
        # Remove punctuation from the word for checking
        clean_word = word.strip(string.punctuation)
        
        # Check if the word is misspelled
        if clean_word and clean_word.lower() not in spell:
            correction = spell.correction(clean_word)
            if correction:  # Only suggest if there's a valid correction
                misspellings[word] = correction

    return misspellings

# Integrate into the main logic
def validate_spelling(input_ppt):
    presentation = Presentation(input_ppt)
    spelling_issues = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip():
                            misspellings = detect_misspellings(run.text)
                            for original_word, correction in misspellings.items():
                                spelling_issues.append({
                                    'slide': slide_index,
                                    'issue': 'Misspelling',
                                    'text': f"Original: {original_word}",
                                    'corrected': f"Suggestion: {correction}"
                                })

    return spelling_issues

# Function to validate fonts in a presentation
def validate_fonts(input_ppt, default_font):
    presentation = Presentation(input_ppt)
    issues = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip():  # Skip empty text
                            # Check for inconsistent fonts
                            if run.font.name != default_font:
                                issues.append({
                                    'slide': slide_index,
                                    'issue': 'Inconsistent Font',
                                    'text': run.text,
                                    'corrected': f"Expected font: {default_font}"
                                })
    return issues

# Function to detect punctuation issues
def validate_punctuation(input_ppt):
    presentation = Presentation(input_ppt)
    punctuation_issues = []

    # Define patterns for punctuation problems
    excessive_punctuation_pattern = r"([!?.:,;]{2,})"  # Two or more punctuation marks
    repeated_word_pattern = r"\b(\w+)\s+\1\b"  # Repeated words (e.g., "the the")

    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip():
                            text = run.text

                            # Check excessive punctuation
                            match = re.search(excessive_punctuation_pattern, text)
                            if match:
                                punctuation_marks = match.group(1)  # Extract punctuation
                                punctuation_issues.append({
                                    'slide': slide_index,
                                    'issue': 'Punctuation Marks',
                                    'text': text,
                                    'corrected': f"Excessive punctuation marks detected ({punctuation_marks})"
                                })

                            # Check repeated words
                            if re.search(repeated_word_pattern, text, flags=re.IGNORECASE):
                                punctuation_issues.append({
                                    'slide': slide_index,
                                    'issue': 'Punctuation Marks',
                                    'text': text,
                                    'corrected': "Repeated words detected"
                                })

    return punctuation_issues

# Function to save issues to CSV
def save_to_csv(issues, output_csv):
    with open(output_csv, mode='w', newline='', encoding='utf-8') as file:
        writer = csv.DictWriter(file, fieldnames=['slide', 'issue', 'text', 'corrected'])
        writer.writeheader()
        writer.writerows(issues)

# Main Streamlit app
def main():
    st.title("PPT Validator")

    uploaded_file = st.file_uploader("Upload a PowerPoint file", type=["pptx"])

    font_options = ["Arial", "Calibri", "Times New Roman", "Verdana", "Helvetica"]
    default_font = st.selectbox("Select the default font for validation", font_options)

    if uploaded_file and st.button("Run Validation"):
        with tempfile.TemporaryDirectory() as tmpdir:
            # Save uploaded file temporarily
            temp_ppt_path = Path(tmpdir) / "uploaded_ppt.pptx"
            with open(temp_ppt_path, "wb") as f:
                f.write(uploaded_file.getbuffer())

            # Output path
            csv_output_path = Path(tmpdir) / "validation_report.csv"

            # Run validations
            font_issues = validate_fonts(temp_ppt_path, default_font)
            punctuation_issues = validate_punctuation(temp_ppt_path)
            grammar_issues = validate_grammar(temp_ppt_path)
            spelling_issues = validate_spelling(temp_ppt_path)

            # Combine issues and save to CSV
            combined_issues = font_issues + punctuation_issues + grammar_issues + spelling_issues
            save_to_csv(combined_issues, csv_output_path)

            # Display success and download link
            st.success("Validation completed!")
            st.download_button("Download Validation Report (CSV)", csv_output_path.read_bytes(),
                               file_name="validation_report.csv")

if __name__ == "__main__":
    main()


########################################################################################################################################

# import streamlit as st
# import tempfile
# from pathlib import Path
# from pptx import Presentation
# from spellchecker import SpellChecker
# import csv
# import re
# import string

# # Function to detect and correct misspellings
# def detect_misspellings(text):
#     spell = SpellChecker()
#     words = text.split()
#     misspellings = {}

#     for word in words:
#         # Remove punctuation from the word for checking
#         clean_word = word.strip(string.punctuation)
        
#         # Check if the word is misspelled
#         if clean_word and clean_word.lower() not in spell:
#             correction = spell.correction(clean_word)
#             if correction:  # Only suggest if there's a valid correction
#                 misspellings[word] = correction

#     return misspellings

# # Integrate into the main logic
# def validate_spelling(input_ppt):
#     presentation = Presentation(input_ppt)
#     spelling_issues = []

#     for slide_index, slide in enumerate(presentation.slides, start=1):
#         for shape in slide.shapes:
#             if shape.has_text_frame:
#                 for paragraph in shape.text_frame.paragraphs:
#                     for run in paragraph.runs:
#                         if run.text.strip():
#                             misspellings = detect_misspellings(run.text)
#                             for original_word, correction in misspellings.items():
#                                 spelling_issues.append({
#                                     'slide': slide_index,
#                                     'issue': 'Misspelling',
#                                     'text': f"Original: {original_word}",
#                                     'corrected': f"Suggestion: {correction}"
#                                 })

#     return spelling_issues

# # Function to validate fonts in a presentation
# def validate_fonts(input_ppt, default_font):
#     presentation = Presentation(input_ppt)
#     issues = []

#     for slide_index, slide in enumerate(presentation.slides, start=1):
#         for shape in slide.shapes:
#             if shape.has_text_frame:
#                 for paragraph in shape.text_frame.paragraphs:
#                     for run in paragraph.runs:
#                         if run.text.strip():  # Skip empty text
#                             # Check for inconsistent fonts
#                             if run.font.name != default_font:
#                                 issues.append({
#                                     'slide': slide_index,
#                                     'issue': 'Inconsistent Font',
#                                     'text': run.text,
#                                     'corrected': f"Expected font: {default_font}"
#                                 })
#     return issues

# # Function to detect punctuation issues
# def validate_punctuation(input_ppt):
#     presentation = Presentation(input_ppt)
#     punctuation_issues = []

#     # Define patterns for punctuation problems
#     excessive_punctuation_pattern = r"([!?.:,;]{2,})"  # Two or more punctuation marks
#     repeated_word_pattern = r"\b(\w+)\s+\1\b"  # Repeated words (e.g., "the the")

#     for slide_index, slide in enumerate(presentation.slides, start=1):
#         for shape in slide.shapes:
#             if shape.has_text_frame:
#                 for paragraph in shape.text_frame.paragraphs:
#                     for run in paragraph.runs:
#                         if run.text.strip():
#                             text = run.text

#                             # Check excessive punctuation
#                             match = re.search(excessive_punctuation_pattern, text)
#                             if match:
#                                 punctuation_marks = match.group(1)  # Extract punctuation
#                                 punctuation_issues.append({
#                                     'slide': slide_index,
#                                     'issue': 'Punctuation Marks',
#                                     'text': text,
#                                     'corrected': f"Excessive punctuation marks detected ({punctuation_marks})"
#                                 })

#                             # Check repeated words
#                             if re.search(repeated_word_pattern, text, flags=re.IGNORECASE):
#                                 punctuation_issues.append({
#                                     'slide': slide_index,
#                                     'issue': 'Punctuation Marks',
#                                     'text': text,
#                                     'corrected': "Repeated words detected"
#                                 })

#     return punctuation_issues

# # Function to save issues to CSV
# def save_to_csv(issues, output_csv):
#     with open(output_csv, mode='w', newline='', encoding='utf-8') as file:
#         writer = csv.DictWriter(file, fieldnames=['slide', 'issue', 'text', 'corrected'])
#         writer.writeheader()
#         writer.writerows(issues)

# # Main Streamlit app
# def main():
#     st.title("PPT Validator")

#     uploaded_file = st.file_uploader("Upload a PowerPoint file", type=["pptx"])

#     font_options = ["Arial", "Calibri", "Times New Roman", "Verdana", "Helvetica"]
#     default_font = st.selectbox("Select the default font for validation", font_options)

#     if uploaded_file and st.button("Run Validation"):
#         with tempfile.TemporaryDirectory() as tmpdir:
#             # Save uploaded file temporarily
#             temp_ppt_path = Path(tmpdir) / "uploaded_ppt.pptx"
#             with open(temp_ppt_path, "wb") as f:
#                 f.write(uploaded_file.getbuffer())

#             # Output path
#             csv_output_path = Path(tmpdir) / "validation_report.csv"

#             # Run validations
#             font_issues = validate_fonts(temp_ppt_path, default_font)
#             punctuation_issues = validate_punctuation(temp_ppt_path)
#             spelling_issues = validate_spelling(temp_ppt_path)

#             # Combine issues and save to CSV
#             combined_issues = font_issues + punctuation_issues + spelling_issues
#             save_to_csv(combined_issues, csv_output_path)

#             # Display success and download link
#             st.success("Validation completed!")
#             st.download_button("Download Validation Report (CSV)", csv_output_path.read_bytes(),
#                                file_name="validation_report.csv")

# if __name__ == "__main__":
#     main()

