import tkinter as tk
from tkinter import ttk, filedialog, messagebox
from pptx import Presentation
from pptx.dml.color import RGBColor
from datetime import datetime
import os
import csv
import re
from pathlib import Path
from matplotlib import font_manager
from transformers import AutoModelForSeq2SeqLM, AutoTokenizer
import logging

# Disable symlink warnings in huggingface_hub
import warnings
os.environ["HF_HUB_DISABLE_SYMLINKS_WARNING"] = "1"
warnings.filterwarnings("ignore", category=UserWarning, module="huggingface_hub")

# Configure logging
logging.basicConfig(
    filename="ppt_validator.log",
    level=logging.DEBUG,
    format="%(asctime)s [%(levelname)s] %(message)s"
)


# Initialize the T5 model for grammar correction
MODEL_NAME = "vennify/t5-base-grammar-correction"
CACHE_DIR = "./model_cache"  # Define cache directory

# Ensure cache directory exists
Path(CACHE_DIR).mkdir(parents=True, exist_ok=True)

# Load or download the model
try:
    tokenizer = AutoTokenizer.from_pretrained(MODEL_NAME, cache_dir=CACHE_DIR)
    model = AutoModelForSeq2SeqLM.from_pretrained(MODEL_NAME, cache_dir=CACHE_DIR)
except Exception as e:
    print(f"Error loading model: {e}. Downloading model because it's missing...")
    tokenizer = AutoTokenizer.from_pretrained(MODEL_NAME, cache_dir=CACHE_DIR, force_download=True)
    model = AutoModelForSeq2SeqLM.from_pretrained(MODEL_NAME, cache_dir=CACHE_DIR, force_download=True)

# Grammar correction function using the T5 model
def correct_grammar(text):
    """
    Correct grammar and spelling using the T5 model.
    """
    try:
        logging.info(f"Original Text: {text}")

        # Split text into chunks if too long for the model
        chunk_size = 200  # Adjust chunk size as per T5's maximum input length
        chunks = [text[i:i + chunk_size] for i in range(0, len(text), chunk_size)]
        corrected_text = []

        for chunk in chunks:
            input_text = f"grammar correction: {chunk}"
            input_ids = tokenizer.encode(input_text, return_tensors="pt")
            outputs = model.generate(input_ids, max_length=512)
            corrected_chunk = tokenizer.decode(outputs[0], skip_special_tokens=True)
            corrected_text.append(corrected_chunk)

        # Join all corrected chunks
        final_text = " ".join(corrected_text)
        logging.info(f"Corrected Text: {final_text}")
        return final_text

    except Exception as e:
        logging.error(f"Error during grammar correction: {e}")
        return text  # Return the original text if correction fails
# Check for inconsistent fonts in the presentation
def check_inconsistent_fonts(file_path, default_font, output_pptx_path=None, progress_callback=None):
    """
    Check for inconsistent fonts in the presentation and highlight them.
    """
    presentation = Presentation(file_path)
    inconsistent_fonts = []
    total_slides = len(presentation.slides)

    for slide_index, slide in enumerate(presentation.slides, start=1):
        if not any(shape.has_text_frame for shape in slide.shapes):
            logging.info(f"Skipping empty slide {slide_index}")
            continue

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        font = getattr(run.font, 'name', None)
                        text = run.text.strip()

                        if font and font != default_font and text:
                            inconsistent_fonts.append({
                                'slide_number': slide_index,
                                'text': text,
                                'font': font
                            })
                            logging.debug(f"Font issue on slide {slide_index}: '{text}' (Font: {font})")
                            try:
                                if run.font and hasattr(run.font, 'color'):
                                    run.font.color.rgb = RGBColor(255, 0, 0)
                            except AttributeError as e:
                                logging.warning(f"Font highlighting error on slide {slide_index}: {e}")

        if progress_callback:
            progress_callback(slide_index / total_slides * 100)

    if output_pptx_path:
        presentation.save(output_pptx_path)

    return inconsistent_fonts

# Check for grammar and spelling issues in the presentation
def check_grammar_and_spelling(file_path, progress_callback=None):
    """
    Use the T5 model to detect grammar and spelling issues in the presentation.
    """
    presentation = Presentation(file_path)
    issues = []
    total_slides = len(presentation.slides)

    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text.extend(
                    run.text.strip() for paragraph in shape.text_frame.paragraphs for run in paragraph.runs
                )

        if not slide_text:
            logging.info(f"Skipping slide {slide_index}: No text found")
            continue

        combined_text = " ".join(slide_text)
        corrected_text = correct_grammar(combined_text)

        def is_significant_change(original, corrected):
            return original.lower().strip() != corrected.lower().strip()

        if is_significant_change(combined_text, corrected_text):
            issues.append({
                'slide_number': slide_index,
                'issue': "Grammar & Spelling Issue",
                'original': combined_text,
                'corrected': corrected_text
            })

        if progress_callback:
            progress_callback(slide_index / total_slides * 100)

    return issues

# Check for punctuation issues in the presentation
def check_punctuation_issues(file_path, progress_callback=None):
    """
    Detect excessive punctuation and repeated words in the presentation.
    """
    presentation = Presentation(file_path)
    punctuation_issues = []
    total_slides = len(presentation.slides)

    for slide_index, slide in enumerate(presentation.slides, start=1):
        if not any(shape.has_text_frame for shape in slide.shapes):
            logging.info(f"Skipping empty slide {slide_index}")
            continue

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text = run.text.strip()
                        if not text:
                            continue

                        # Regex pattern for excessive punctuation
                        excessive_punctuation_pattern = r"(?:[!?.:,;]{2,})"  # Match two or more punctuation marks
                        repeated_word_pattern = r"\b(\w+)\s+\1\b"  # Match repeated words like "the the"

                        # Find excessive punctuation
                        excessive_punctuation_matches = re.findall(excessive_punctuation_pattern, text)
                        for match in excessive_punctuation_matches:
                            punctuation_issues.append({
                                'slide_number': slide_index,
                                'text': text,
                                'issue': f"Excessive punctuation: {match}"
                            })
                            logging.debug(f"Excessive punctuation found on slide {slide_index}: {text}")

                        # Find repeated words
                        repeated_word_matches = re.findall(repeated_word_pattern, text, flags=re.IGNORECASE)
                        for match in repeated_word_matches:
                            punctuation_issues.append({
                                'slide_number': slide_index,
                                'text': text,
                                'issue': f"Repeated word: {match}"
                            })
                            logging.debug(f"Repeated word found on slide {slide_index}: {text}")

        if progress_callback:
            progress_callback(slide_index / total_slides * 100)

    return punctuation_issues

# Save detected issues to a CSV file
def save_issues_to_csv(font_issues, grammar_and_spelling_issues, punctuation_issues, output_name):
    """
    Save detected issues to a CSV file.
    """
    downloads_path = Path.home() / "Downloads"
    csv_file_name = downloads_path / f"{output_name}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv"

    try:
        with open(csv_file_name, mode='w', newline='', encoding='utf-8-sig') as file:
            writer = csv.writer(file)
            writer.writerow(["Slide Number", "Issue Type", "Original Text", "Corrected Text or Details"])

            for issue in font_issues:
                writer.writerow([
                    issue['slide_number'],
                    "Inconsistent Font Issue",
                    issue['text'],
                    f"Font: {issue['font']} (expected: {default_font_var.get()})"
                ])

            for issue in grammar_and_spelling_issues:
                writer.writerow([
                    issue['slide_number'],
                    issue['issue'],
                    issue['original'],
                    issue['corrected']
                ])

            for issue in punctuation_issues:
                writer.writerow([
                    issue['slide_number'],
                    "Punctuation Issue",
                    issue['text'],
                    issue['issue']
                ])

        logging.info(f"CSV saved successfully: {csv_file_name}")
        return csv_file_name

    except Exception as e:
        logging.error(f"Error saving CSV file: {e}")
        return None  # Return None if saving fails

# GUI Components
root = tk.Tk()
root.title("PowerPoint Validator")
root.geometry("600x450")  # Set window size

# Frame for file selection
frame_file = tk.Frame(root)
frame_file.pack(pady=10)
tk.Label(frame_file, text="Select PowerPoint File:").grid(row=0, column=0, padx=5, pady=5)

# Entry box to display the file path
entry_file_path = tk.Entry(frame_file, width=40)
entry_file_path.grid(row=0, column=1, padx=5, pady=5)

# Function to browse for a file
def select_file():
    file_path = filedialog.askopenfilename(filetypes=[("PowerPoint files", "*.pptx"), ("All files", "*.*")])
    if file_path:
        entry_file_path.delete(0, tk.END)  # Clear the entry box
        entry_file_path.insert(0, file_path)  # Insert selected file path

tk.Button(frame_file, text="Browse", command=select_file).grid(row=0, column=2, padx=5, pady=5)

# Dropdown menu for selecting the default font
tk.Label(root, text="Select Default Font:").pack(pady=5)
default_font_var = tk.StringVar(root)
fonts = sorted(set(f.name for f in font_manager.fontManager.ttflist))  # Get available fonts
default_font_dropdown = ttk.Combobox(root, textvariable=default_font_var, values=fonts, width=40)
default_font_dropdown.pack(pady=5)

# Progress bar and progress label
progress_label = tk.Label(root, text="")
progress_label.pack(pady=5)
progress_bar = ttk.Progressbar(root, orient="horizontal", length=400, mode="determinate")
progress_bar.pack(pady=5)
progress_percentage = tk.Label(root, text="0%")
progress_percentage.pack(pady=5)

# Function to validate and process the PowerPoint file
def run_tool():
    file_path = entry_file_path.get()  # Get file path from the entry box
    default_font = default_font_var.get()  # Get selected font

    # Validate file extension
    if not file_path.endswith('.pptx'):
        messagebox.showerror("Error", "Please select a valid PowerPoint (.pptx) file.")
        logging.warning("Invalid file extension selected.")
        return

    # Validate file existence
    if not os.path.exists(file_path):
        messagebox.showerror("Error", f"PowerPoint file not found: {file_path}")
        logging.error("PowerPoint file not found.")
        return

    # Validate default font selection
    if not default_font:
        messagebox.showerror("Error", "Default font not selected. Please select a font from the dropdown.")
        logging.error("Default font not selected.")
        return

    # Set output paths
    downloads_path = Path.home() / "Downloads"
    output_pptx_path = downloads_path / (Path(file_path).stem + "_highlighted.pptx")

    # Function to update the progress bar
    def update_progress(value):
        progress_bar['value'] = value
        progress_percentage.config(text=f"{int(value)}%")
        root.update_idletasks()

    try:
        logging.info("Starting validation...")

        # Step 1: Check for inconsistent fonts
        progress_label.config(text="Checking Inconsistent Fonts...")
        font_issues = check_inconsistent_fonts(file_path, default_font, output_pptx_path, update_progress)

        # Step 2: Check for grammar and spelling issues
        progress_label.config(text="Checking Grammar and Spelling...")
        grammar_issues = check_grammar_and_spelling(file_path, update_progress)

        # Step 3: Check for punctuation issues
        progress_label.config(text="Checking Punctuation Issues...")
        punctuation_issues = check_punctuation_issues(file_path, update_progress)

        # Step 4: Save the results to a CSV file
        progress_label.config(text="Saving Results to CSV...")
        csv_file = save_issues_to_csv(font_issues, grammar_issues, punctuation_issues, "Validation_Report")

        # Final update of progress
        update_progress(100)
        if csv_file:
            messagebox.showinfo(
                "Success",
                f"Validation completed.\n"
                f"CSV Report: {csv_file}\n"
                f"Highlighted PPTX: {output_pptx_path}"
            )
            logging.info("Validation completed successfully.")
        else:
            # Notify user when no issues are found
            messagebox.showinfo(
                "Validation Completed",
                "Validation completed. No issues found in the slides."
            )
            logging.info(f"No issues found for file: {file_path}")

        progress_label.config(text="Validation Completed.")

    except Exception as e:
        messagebox.showerror("Error", f"An unexpected error occurred: {e}")
        logging.error(f"Unexpected error: {e}")


# Button to start the validation process
tk.Button(root, text="Run Validation", command=run_tool).pack(pady=20)

# Start the GUI event loop
root.mainloop()